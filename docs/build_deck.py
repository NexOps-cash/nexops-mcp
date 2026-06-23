"""
NexOps Protocol — Pitch Deck v2
Fixes: no overclaiming, accurate RCAs, gray cards / white text,
       green as accent-only (~10%), added What We Learned slide,
       cover value-prop, story-format Slide 9, clean tier matrix Slide 4.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import math

# ─── Palette ──────────────────────────────────────────────
# Rule: green is accent only. Cards are dark gray. Titles are white.
C_BG      = RGBColor(0x07, 0x0c, 0x10)   # slide background – near-black
C_CARD    = RGBColor(0x11, 0x18, 0x27)   # primary card (gray-900)
C_CARD2   = RGBColor(0x1f, 0x29, 0x37)   # secondary card (gray-800)
C_GREEN   = RGBColor(0x22, 0xc5, 0x5e)   # accent green  – use sparingly
C_GREEN_D = RGBColor(0x16, 0x7a, 0x3b)   # dark green   – accent fills
C_GREEN_L = RGBColor(0x4a, 0xde, 0x80)   # light green
C_YELLOW  = RGBColor(0xfb, 0xbf, 0x24)   # amber / warning
C_ORANGE  = RGBColor(0xf9, 0x73, 0x16)   # orange / in-progress
C_CYAN    = RGBColor(0x06, 0xb6, 0xd4)   # cyan – rare accent
C_WHITE   = RGBColor(0xff, 0xff, 0xff)
C_LGRAY   = RGBColor(0xd1, 0xd5, 0xdb)   # light gray – body text
C_GRAY    = RGBColor(0x9c, 0xa3, 0xaf)   # mid gray – captions

W = Inches(13.33)
H = Inches(7.5)

# ─── Primitives ───────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def oval(slide, l, t, w, h, fill, line=None, lw=Pt(1)):
    s = slide.shapes.add_shape(9, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = lw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=Pt(13), bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def multiline(slide, lines, l, t, w, h,
              default_size=Pt(13), default_bold=False,
              default_color=C_WHITE, default_align=PP_ALIGN.LEFT):
    """lines = list of str or (text, size, bold, color, align)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            tx, sz, bld, clr, al = item, default_size, default_bold, default_color, default_align
        else:
            tx = item[0]
            sz  = item[1] if len(item) > 1 else default_size
            bld = item[2] if len(item) > 2 else default_bold
            clr = item[3] if len(item) > 3 else default_color
            al  = item[4] if len(item) > 4 else default_align
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = al
        r = p.add_run(); r.text = tx
        r.font.size = sz; r.font.bold = bld; r.font.color.rgb = clr
    return tb

# ─── Layout helpers ───────────────────────────────────────

def bottom_bar(slide):
    """1 px green accent line at bottom"""
    rect(slide, 0, H - Inches(0.055), W, Inches(0.055), C_GREEN)

def slide_num(slide, n):
    txt(slide, f"{n:02d}", W - Inches(0.85), Inches(0.14),
        Inches(0.65), Inches(0.3), size=Pt(10), color=C_GREEN,
        align=PP_ALIGN.RIGHT)
    bottom_bar(slide)

def tag(slide, label, x=Inches(0.55), y=Inches(0.17)):
    """Small pill label in top-left"""
    rect(slide, x, y, Inches(2.0), Inches(0.27), C_GREEN_D)
    txt(slide, label, x + Inches(0.08), y + Inches(0.02),
        Inches(1.85), Inches(0.25), size=Pt(8), bold=True, color=C_GREEN)

def title(slide, text, y=Inches(0.56), size=Pt(32)):
    txt(slide, text, Inches(0.55), y, W - Inches(1.1), Inches(0.65),
        size=size, bold=True, color=C_WHITE)

def subtitle_line(slide, text, y=Inches(1.28)):
    txt(slide, text, Inches(0.55), y, W - Inches(1.1), Inches(0.35),
        size=Pt(12), color=C_GRAY)

def divider(slide, y, x=Inches(0.55), w=None, color=C_CARD2):
    if w is None: w = W - Inches(1.1)
    rect(slide, x, y, w, Inches(0.025), color)

def accent_bar_left(slide, y, h, color=C_GREEN):
    """Thin vertical green line on left of a card"""
    rect(slide, Inches(0.55), y, Inches(0.055), h, color)

# ─────────────────────────────────────────────────────────
# SLIDE 1 — COVER  (value prop added, green as accent only)
# ─────────────────────────────────────────────────────────
def slide_01(prs):
    s = blank(prs); bg(s)

    # Subtle dot grid
    for row in range(8):
        for col in range(20):
            x = Inches(0.15 + col * 0.67)
            y = Inches(0.3 + row * 0.9 + (0.45 if col % 2 else 0))
            o = oval(s, x, y, Inches(0.05), Inches(0.05),
                     RGBColor(0x1e, 0x29, 0x3b))

    # Green accent line left
    rect(s, Inches(0.55), Inches(1.25), Inches(0.06), Inches(2.1), C_GREEN)

    # Brand
    txt(s, "NEXOPS", Inches(0.75), Inches(1.2), Inches(7), Inches(1.1),
        size=Pt(68), bold=True, color=C_WHITE)
    txt(s, "PROTOCOL", Inches(0.76), Inches(2.2), Inches(7), Inches(0.7),
        size=Pt(42), bold=False, color=C_GREEN)

    # Headline
    txt(s, "Security-First Smart Contract Infrastructure for Bitcoin Cash",
        Inches(0.76), Inches(3.05), Inches(9.5), Inches(0.52),
        size=Pt(17), color=C_LGRAY)

    # Value proposition (new)
    rect(s, Inches(0.55), Inches(3.75), Inches(9.8), Inches(0.52), C_CARD2)
    rect(s, Inches(0.55), Inches(3.75), Inches(0.055), Inches(0.52), C_GREEN)
    txt(s,
        "Generate, verify, and benchmark BCH smart contracts before deployment.",
        Inches(0.72), Inches(3.78), Inches(9.5), Inches(0.45),
        size=Pt(13), bold=True, color=C_WHITE)

    # Four pillars — white text, green bottom line only
    pillars = ["Generation", "Verification", "Benchmarking", "Audit Intelligence"]
    for i, p in enumerate(pillars):
        px = Inches(0.76) + i * Inches(2.4)
        rect(s, px, Inches(4.52), Inches(2.2), Inches(0.55), C_CARD)
        rect(s, px, Inches(5.05), Inches(2.2), Inches(0.03), C_GREEN)
        txt(s, p, px + Inches(0.12), Inches(4.56),
            Inches(2.0), Inches(0.42), size=Pt(12), bold=True, color=C_WHITE)

    # Right network graphic — minimal
    cx, cy = Inches(11.35), Inches(3.5)
    # Center node
    oval(s, cx - Inches(0.42), cy - Inches(0.42), Inches(0.84), Inches(0.84),
         C_GREEN)
    txt(s, "NX", cx - Inches(0.42), cy - Inches(0.2),
        Inches(0.84), Inches(0.35), size=Pt(12), bold=True,
        color=C_BG, align=PP_ALIGN.CENTER)

    nodes = [("BCH", C_YELLOW), ("FT", C_GREEN_L), ("NFT", C_CYAN),
             ("Vault", C_LGRAY), ("Audit", C_ORANGE), ("Bench", C_GRAY)]
    for i, (lbl, clr) in enumerate(nodes):
        ang = i * math.pi * 2 / 6 - math.pi / 2
        r = Inches(1.55)
        nx = cx + r * math.cos(ang) - Inches(0.28)
        ny = cy + r * math.sin(ang) - Inches(0.28)
        # connector line
        lx = cx + (r - Inches(0.35)) * math.cos(ang)
        ly = cy + (r - Inches(0.35)) * math.sin(ang)
        rect(s, cx - Inches(0.01), cy - Inches(0.01),
             Inches(0.02), Inches(0.02), clr)  # tiny dot
        oval(s, nx, ny, Inches(0.56), Inches(0.56), C_CARD, clr, Pt(1.2))
        txt(s, lbl, nx, ny + Inches(0.12), Inches(0.56), Inches(0.28),
            size=Pt(8), bold=True, color=clr, align=PP_ALIGN.CENTER)

    # Bottom strip
    rect(s, 0, H - Inches(0.55), W, Inches(0.5), C_CARD)
    txt(s, "BCH Hackcelerator  •  Security Infrastructure  •  2026",
        Inches(0.55), H - Inches(0.48), W - Inches(1.1), Inches(0.38),
        size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)
    bottom_bar(s)


# ─────────────────────────────────────────────────────────
# SLIDE 2 — WHERE WE STARTED
# ─────────────────────────────────────────────────────────
def slide_02(prs):
    s = blank(prs); bg(s)
    slide_num(s, 2); tag(s, "ORIGIN STORY")
    title(s, "Where We Started")
    subtitle_line(s, "Original BCH Hackcelerator scope — pattern-support MVP")
    divider(s, Inches(1.65))

    # Col 1 – Patterns
    rect(s, Inches(0.55), Inches(1.75), Inches(3.85), Inches(3.9), C_CARD)
    rect(s, Inches(0.55), Inches(1.75), Inches(3.85), Inches(0.42), C_CARD2)
    txt(s, "PATTERN COVERAGE", Inches(0.7), Inches(1.78),
        Inches(3.5), Inches(0.35), size=Pt(9), bold=True, color=C_GREEN)

    patterns = ["Single Signature Transfer", "Timelock", "Hashlock",
                "Multisig", "Escrow", "Refundable Payment",
                "Split Payment", "Vault", "Covenant",
                "Conditional Spend", "Decay"]
    for i, p in enumerate(patterns):
        rect(s, Inches(0.72), Inches(2.32) + i * Inches(0.3),
             Inches(0.12), Inches(0.12), C_GREEN)
        txt(s, p, Inches(0.95), Inches(2.26) + i * Inches(0.3),
            Inches(3.2), Inches(0.28), size=Pt(10), color=C_LGRAY)

    # Col 2 – CashTokens
    rect(s, Inches(4.6), Inches(1.75), Inches(3.5), Inches(2.4), C_CARD)
    rect(s, Inches(4.6), Inches(1.75), Inches(3.5), Inches(0.42), C_CARD2)
    txt(s, "CASHTOKENS", Inches(4.75), Inches(1.78),
        Inches(3.2), Inches(0.35), size=Pt(9), bold=True, color=C_CYAN)
    ct = ["FT Mint", "NFT Mint & Transfer", "Category Validation",
          "Amount Validation", "Invalid Token Detection"]
    for i, c in enumerate(ct):
        rect(s, Inches(4.77), Inches(2.32) + i * Inches(0.35),
             Inches(0.12), Inches(0.12), C_CYAN)
        txt(s, c, Inches(5.0), Inches(2.26) + i * Inches(0.35),
            Inches(2.9), Inches(0.28), size=Pt(10), color=C_LGRAY)

    # Arrow
    txt(s, ">", Inches(8.3), Inches(3.0), Inches(0.5), Inches(0.5),
        size=Pt(32), bold=True, color=C_GREEN_D, align=PP_ALIGN.CENTER)

    # "Became" box
    rect(s, Inches(9.0), Inches(1.75), Inches(3.78), Inches(3.6), C_CARD)
    rect(s, Inches(9.0), Inches(1.75), Inches(0.055), Inches(3.6), C_GREEN)
    txt(s, "WHAT IT BECAME", Inches(9.18), Inches(1.82),
        Inches(3.4), Inches(0.35), size=Pt(9), bold=True, color=C_GREEN)

    became = [
        ("Security Validation Engine", C_GREEN_L),
        ("Benchmark Infrastructure",   C_GREEN),
        ("Semantic Audit System",       C_YELLOW),
        ("CashTokens Full Suite",       C_CYAN),
        ("Contract Intelligence Platform", C_ORANGE),
    ]
    for i, (m, c) in enumerate(became):
        rect(s, Inches(9.18), Inches(2.35) + i * Inches(0.52),
             Inches(0.12), Inches(0.12), c)
        txt(s, m, Inches(9.42), Inches(2.3) + i * Inches(0.52),
            Inches(3.22), Inches(0.35), size=Pt(11), color=C_WHITE)

    # Key message
    rect(s, Inches(0.55), Inches(5.85), W - Inches(1.1), Inches(0.48), C_CARD2)
    rect(s, Inches(0.55), Inches(5.85), Inches(0.055), Inches(0.48), C_GREEN)
    txt(s, "Started as a pattern-support MVP — evolved into a full security infrastructure platform.",
        Inches(0.72), Inches(5.89), W - Inches(1.3), Inches(0.38),
        size=Pt(11), bold=True, color=C_WHITE)


# ─────────────────────────────────────────────────────────
# SLIDE 3 — EVOLUTION TIMELINE
# ─────────────────────────────────────────────────────────
def slide_03(prs):
    s = blank(prs); bg(s)
    slide_num(s, 3); tag(s, "EVOLUTION")
    title(s, "What NexOps Became")
    subtitle_line(s, "From MVP generator to security infrastructure platform")
    divider(s, Inches(1.65))

    stages = [
        ("01", "MVP\nGenerator",        "BCH pattern generation",           C_GRAY,    Inches(0.55)),
        ("02", "Security\nValidation",  "Structural & semantic checks",      C_GREEN_L, Inches(3.1)),
        ("03", "Benchmark\nInfrastructure","Repeatable evaluation suites",   C_GREEN,   Inches(5.65)),
        ("04", "Semantic\nAudit",       "Intent-aware finding policies",     C_YELLOW,  Inches(8.2)),
        ("05", "Contract\nIntelligence","Composition & full platform",       C_ORANGE,  Inches(10.75)),
    ]

    # Spine
    rect(s, Inches(0.55), Inches(3.6), Inches(12.23), Inches(0.04), C_CARD2)

    for num, ttl, sub, clr, x in stages:
        node_y = Inches(3.4)
        oval(s, x, node_y, Inches(0.44), Inches(0.44), C_CARD, clr, Pt(2))
        txt(s, num, x, node_y, Inches(0.44), Inches(0.4),
            size=Pt(9), bold=True, color=clr, align=PP_ALIGN.CENTER)

        # Card below
        rect(s, x - Inches(0.28), Inches(4.1), Inches(2.3), Inches(2.2), C_CARD)
        rect(s, x - Inches(0.28), Inches(4.1), Inches(0.055), Inches(2.2), clr)
        txt(s, ttl, x, Inches(4.18), Inches(2.0), Inches(0.62),
            size=Pt(13), bold=True, color=C_WHITE)
        txt(s, sub, x, Inches(4.88), Inches(2.0), Inches(0.65),
            size=Pt(10), color=C_GRAY)
        rect(s, x + Inches(0.18), Inches(3.82), Inches(0.03), Inches(0.3), clr)

    for i in range(4):
        ax = Inches(1.25) + i * Inches(2.55)
        txt(s, ">", ax, Inches(3.38), Inches(0.5), Inches(0.38),
            size=Pt(14), color=C_CARD2, align=PP_ALIGN.CENTER)

    rect(s, Inches(0.55), Inches(6.42), W - Inches(1.1), Inches(0.52), C_CARD2)
    rect(s, Inches(0.55), Inches(6.42), Inches(0.055), Inches(0.52), C_GREEN)
    txt(s,
        "Scope expanded far beyond contract generation — NexOps now owns the full contract lifecycle.",
        Inches(0.72), Inches(6.46), W - Inches(1.3), Inches(0.42),
        size=Pt(11), bold=True, color=C_WHITE)


# ─────────────────────────────────────────────────────────
# SLIDE 4 — PATTERN COVERAGE (tier-based, no %)
# ─────────────────────────────────────────────────────────
def slide_04(prs):
    s = blank(prs); bg(s)
    slide_num(s, 4); tag(s, "PATTERN COVERAGE")
    title(s, "Contract Readiness — Honest Assessment")
    subtitle_line(s, "Status based on benchmark convergence, not theoretical coverage")
    divider(s, Inches(1.65))

    tiers = [
        ("Production Converged", C_GREEN,  C_GREEN_D,
         ["Escrow", "Multisig", "Timelock", "Hashlock"],
         "Positive benchmark paths reach 1.0 score. Stable across evaluator runs."),
        ("Advanced Coverage",    C_YELLOW, RGBColor(0x3d, 0x2f, 0x05),
         ["Vault", "Conditional Spend", "Covenant"],
         "Generation solid. Audit + benchmark coverage growing. Active development."),
        ("In Progress",          C_ORANGE, RGBColor(0x3d, 0x18, 0x03),
         ["Split Payment", "Refundable Payment"],
         "Known composition and routing blockers under active investigation."),
    ]

    tw = Inches(3.7)
    for ti, (tier_name, clr, bg_clr, items, note) in enumerate(tiers):
        tx = Inches(0.55) + ti * (tw + Inches(0.22))
        rect(s, tx, Inches(1.82), tw, Inches(4.8), C_CARD)
        rect(s, tx, Inches(1.82), tw, Inches(0.55), bg_clr)

        # Tier dot + name
        oval(s, tx + Inches(0.14), Inches(1.96), Inches(0.22), Inches(0.22), clr)
        txt(s, tier_name, tx + Inches(0.46), Inches(1.9),
            tw - Inches(0.55), Inches(0.42), size=Pt(12), bold=True, color=C_WHITE)

        # Pattern list
        for ii, pat in enumerate(items):
            iy = Inches(2.56) + ii * Inches(0.62)
            rect(s, tx + Inches(0.2), iy + Inches(0.15),
                 Inches(0.08), Inches(0.28), clr)
            txt(s, pat, tx + Inches(0.42), iy + Inches(0.08),
                tw - Inches(0.55), Inches(0.42), size=Pt(13), color=C_WHITE)

        # Note at bottom
        rect(s, tx, Inches(5.58), tw, Inches(1.0), RGBColor(0x0d, 0x13, 0x1f))
        txt(s, note, tx + Inches(0.14), Inches(5.65),
            tw - Inches(0.25), Inches(0.9), size=Pt(10), color=C_GRAY)

    # Clarification footnote
    txt(s,
        '"Production Converged" = positive benchmark paths reach score 1.0 across evaluator runs.',
        Inches(0.55), Inches(6.88), W - Inches(1.1), Inches(0.32),
        size=Pt(9), color=C_GRAY, italic=True)


# ─────────────────────────────────────────────────────────
# SLIDE 5 — CASHTOKENS EXPANSION
# ─────────────────────────────────────────────────────────
def slide_05(prs):
    s = blank(prs); bg(s)
    slide_num(s, 5); tag(s, "CASHTOKENS")
    title(s, "CashTokens Full Ecosystem Support")
    subtitle_line(s, "From basic FT minting to hybrid covenant-token systems")
    divider(s, Inches(1.65))

    columns = [
        ("Fungible Tokens",   C_GREEN, [
            ("Minting & Issuance",   "Token creation with genesis rules"),
            ("Supply Enforcement",    "Immutable cap verification"),
            ("Transfer Validation",   "Balance conservation checks"),
            ("Balance Assertions",    "On-chain supply invariants"),
        ]),
        ("Non-Fungible Tokens", C_CYAN, [
            ("NFT Minting",          "Unique category issuance"),
            ("Transfer Logic",       "Ownership transfer rules"),
            ("Mutable State Updates","Commitment field mutation"),
            ("Immutable NFTs",       "Lock-after-mint enforcement"),
        ]),
        ("Advanced Patterns",  C_YELLOW, [
            ("Authority Custody",    "Controlled minting keys"),
            ("Category Continuity",  "Cross-tx category preservation"),
            ("Amount Preservation",  "Token conservation covenants"),
            ("Hybrid Covenants",     "Script + token logic combined"),
        ]),
    ]

    cw = Inches(3.85)
    for ci, (col_title, clr, items) in enumerate(columns):
        cx = Inches(0.55) + ci * (cw + Inches(0.2))
        rect(s, cx, Inches(1.82), cw, Inches(4.6), C_CARD)
        rect(s, cx, Inches(1.82), cw, Inches(0.48), C_CARD2)
        rect(s, cx, Inches(1.82), Inches(0.055), Inches(4.6), clr)
        txt(s, col_title, cx + Inches(0.2), Inches(1.86),
            cw - Inches(0.28), Inches(0.38), size=Pt(11), bold=True, color=clr)

        for ii, (item_title, item_desc) in enumerate(items):
            iy = Inches(2.48) + ii * Inches(0.95)
            rect(s, cx + Inches(0.2), iy + Inches(0.08),
                 Inches(0.1), Inches(0.3), clr)
            txt(s, item_title, cx + Inches(0.42), iy + Inches(0.04),
                cw - Inches(0.55), Inches(0.35), size=Pt(12), bold=True, color=C_WHITE)
            txt(s, item_desc, cx + Inches(0.42), iy + Inches(0.42),
                cw - Inches(0.55), Inches(0.35), size=Pt(10), color=C_GRAY)

    # Benchmark proof points strip
    rect(s, Inches(0.55), Inches(6.55), W - Inches(1.1), Inches(0.68), C_CARD2)
    rect(s, Inches(0.55), Inches(6.55), Inches(0.055), Inches(0.68), C_CYAN)
    txt(s, "10 Executable CashTokens Benchmarks  —  Dedicated benchmark suite per token category",
        Inches(0.72), Inches(6.58), Inches(7.5), Inches(0.35),
        size=Pt(11), bold=True, color=C_WHITE)
    proof_tags = ["Authority leak coverage", "Token inflation coverage",
                  "NFT integrity coverage", "Hybrid continuity coverage"]
    for ti, tag_text in enumerate(proof_tags):
        tx = Inches(0.72) + ti * Inches(2.98)
        rect(s, tx, Inches(6.95), Inches(2.75), Inches(0.22), C_CARD)
        txt(s, tag_text, tx + Inches(0.08), Inches(6.96), Inches(2.6), Inches(0.2),
            size=Pt(9), color=C_CYAN)


# ─────────────────────────────────────────────────────────
# SLIDE 6 — SECURITY ENGINE
# ─────────────────────────────────────────────────────────
def slide_06(prs):
    s = blank(prs); bg(s)
    slide_num(s, 6); tag(s, "SECURITY ENGINE")
    title(s, "Multi-Layer Security Architecture")
    subtitle_line(s, "Every contract passes through a structured validation pipeline")
    divider(s, Inches(1.65))

    pipeline = [
        ("01", "User Intent",           "Parameters or natural language",    C_GRAY),
        ("02", "Generator",             "CashScript template engine",        C_GREEN_L),
        ("03", "Structural Integrity",  "Schema & constraint checks",        C_GREEN),
        ("04", "Semantic Validation",   "Intent invariant verification",     C_YELLOW),
        ("05", "Compile Verification",  "CashScript compile pass",           C_CYAN),
        ("06", "Security Judge",        "Capability & finding policies",     C_ORANGE),
        ("07", "Findings",              "Structured, actionable audit report",C_GREEN),
    ]

    sw = Inches(5.7)
    sh = Inches(0.68)
    sx = Inches(0.55)

    for i, (num, label, sub, clr) in enumerate(pipeline):
        sy = Inches(1.82) + i * (sh + Inches(0.05))
        rect(s, sx, sy, sw, sh, C_CARD)
        rect(s, sx, sy, Inches(0.055), sh, clr)
        txt(s, num, sx + Inches(0.14), sy + Inches(0.18),
            Inches(0.38), Inches(0.35), size=Pt(9), bold=True, color=clr)
        txt(s, label, sx + Inches(0.62), sy + Inches(0.1),
            Inches(2.8), Inches(0.42), size=Pt(13), bold=True, color=C_WHITE)
        txt(s, sub, sx + Inches(3.5), sy + Inches(0.12),
            Inches(2.05), Inches(0.44), size=Pt(9), color=C_GRAY)
        if i < 6:
            txt(s, "v", sx + Inches(2.6), sy + sh - Inches(0.06),
                Inches(0.4), Inches(0.22), size=Pt(9), color=clr,
                align=PP_ALIGN.CENTER)

    # Right panel
    rx = Inches(6.7)
    txt(s, "KEY CAPABILITIES", rx, Inches(1.82),
        Inches(6.08), Inches(0.35), size=Pt(9), bold=True, color=C_GREEN)
    divider(s, Inches(2.22), x=rx, w=Inches(6.08))

    caps = [
        ("Capability Extraction",    C_GREEN,  "Parses all actions a contract can perform"),
        ("Intent Invariants",        C_YELLOW, "Validates what the contract should do"),
        ("Semantic Security Judge",  C_ORANGE, "Detects dangerous capability mismatches"),
        ("Finding Policies",         C_CYAN,   "Structured, machine-readable output"),
    ]
    for i, (cap, clr, desc) in enumerate(caps):
        cy = Inches(2.32) + i * Inches(1.2)
        rect(s, rx, cy, Inches(6.08), Inches(1.05), C_CARD)
        rect(s, rx, cy, Inches(0.055), Inches(1.05), clr)
        txt(s, cap, rx + Inches(0.18), cy + Inches(0.1),
            Inches(5.75), Inches(0.42), size=Pt(13), bold=True, color=C_WHITE)
        txt(s, desc, rx + Inches(0.18), cy + Inches(0.58),
            Inches(5.75), Inches(0.38), size=Pt(11), color=C_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE 7 — SECURITY AUDIT EVOLUTION (new)
# ─────────────────────────────────────────────────────────
def slide_audit_evolution(prs):
    s = blank(prs); bg(s)
    slide_num(s, 7); tag(s, "AUDIT EVOLUTION")
    title(s, "Security Audit — How We Got Here")
    subtitle_line(s, "Four generations of audit thinking — each built on the failures of the previous")
    divider(s, Inches(1.65))

    versions = [
        ("V1",    "LLM Findings",
         "Unstructured",
         ["Natural language output", "No reproducibility", "No baseline comparison", "Results varied per run"],
         C_GRAY),
        ("V2",    "Intent Invariants",
         "Structured",
         ["Defined what contract SHOULD do", "Invariant specification layer", "Deterministic output", "First structured audit format"],
         C_CYAN),
        ("V2.1",  "Semantic Security Judge",
         "Capability-aware",
         ["Parsed what contract CAN do", "Compared intent vs capability", "Mismatch = finding", "Actionable finding policies"],
         C_YELLOW),
        ("TODAY", "Replay + Benchmark\n+ Adversarial",
         "Validated",
         ["33 replay regression cases", "200 adversarial scenarios", "Evaluator corpus tested", "Results reproducible + diffable"],
         C_GREEN),
    ]

    vw = Inches(2.92)
    for vi, (ver, title_v, status, bullets, clr) in enumerate(versions):
        vx = Inches(0.55) + vi * (vw + Inches(0.2))

        # Column card
        rect(s, vx, Inches(1.82), vw, Inches(5.12), C_CARD)
        rect(s, vx, Inches(1.82), vw, Inches(0.06), clr)

        # Version badge
        rect(s, vx + Inches(0.15), Inches(1.98), Inches(0.72), Inches(0.38), C_CARD2)
        txt(s, ver, vx + Inches(0.15), Inches(2.0), Inches(0.72), Inches(0.3),
            size=Pt(10), bold=True, color=clr, align=PP_ALIGN.CENTER)

        # Status pill
        txt(s, status, vx + Inches(0.98), Inches(2.02), vw - Inches(1.1), Inches(0.28),
            size=Pt(9), color=C_GRAY)

        # Title
        txt(s, title_v, vx + Inches(0.15), Inches(2.52), vw - Inches(0.28), Inches(0.65),
            size=Pt(13), bold=True, color=C_WHITE)

        # Divider
        rect(s, vx + Inches(0.15), Inches(3.25), vw - Inches(0.28), Inches(0.025), C_CARD2)

        # Bullets
        for bi, bullet in enumerate(bullets):
            by = Inches(3.38) + bi * Inches(0.58)
            rect(s, vx + Inches(0.2), by + Inches(0.1), Inches(0.07), Inches(0.22), clr)
            txt(s, bullet, vx + Inches(0.38), by + Inches(0.04),
                vw - Inches(0.5), Inches(0.42), size=Pt(10), color=C_LGRAY)

        # Arrows between columns
        if vi < 3:
            txt(s, ">", vx + vw + Inches(0.0), Inches(3.4),
                Inches(0.22), Inches(0.38), size=Pt(14), color=C_CARD2,
                align=PP_ALIGN.CENTER)

    rect(s, Inches(0.55), Inches(7.08), W - Inches(1.1), Inches(0.28), C_CARD2)
    txt(s, "Each version was driven by a real failure mode in the previous — not roadmap planning.",
        Inches(0.72), Inches(7.1), W - Inches(1.3), Inches(0.24),
        size=Pt(9), color=C_GRAY, italic=True)


# ─────────────────────────────────────────────────────────
# SLIDE 8 — BENCHMARKING SYSTEM
# ─────────────────────────────────────────────────────────
def slide_07(prs):
    s = blank(prs); bg(s)
    slide_num(s, 8); tag(s, "BENCHMARK SYSTEM")
    title(s, "Repeatable Evaluation Infrastructure")
    subtitle_line(s, "Moved from ad-hoc manual testing to structured, versioned benchmark suites")
    divider(s, Inches(1.65))

    for bi, (label, items, clr, bx) in enumerate([
        ("BEFORE — Manual Validation", [
            ("Ad-hoc test scripts per pattern",     False),
            ("No baseline or regression tracking",  False),
            ("Results not reproducible",            False),
            ("Coverage unknown — no metrics",       False),
            ("Bugs indistinguishable from gaps",    False),
        ], C_ORANGE, Inches(0.55)),
        ("AFTER — NexOps Benchmarks", [
            ("Executable benchmark suites",         True),
            ("Pattern + CashTokens coverage",       True),
            ("Adversarial edge-case evaluations",   True),
            ("Regression testing built-in",         True),
            ("Quantified readiness per pattern",    True),
        ], C_GREEN, Inches(7.0)),
    ]):
        cw = Inches(5.7)
        rect(s, bx, Inches(1.82), cw, Inches(4.65), C_CARD)
        rect(s, bx, Inches(1.82), cw, Inches(0.48), C_CARD2)
        rect(s, bx, Inches(1.82), Inches(0.055), Inches(4.65), clr)
        txt(s, label, bx + Inches(0.2), Inches(1.86),
            cw - Inches(0.28), Inches(0.38), size=Pt(10), bold=True, color=clr)
        for ii, (item, good) in enumerate(items):
            iy = Inches(2.48) + ii * Inches(0.68)
            ic = C_GREEN if good else C_ORANGE
            mark = "+" if good else "-"
            rect(s, bx + Inches(0.22), iy + Inches(0.1),
                 Inches(0.3), Inches(0.3), ic)
            txt(s, mark, bx + Inches(0.22), iy + Inches(0.06),
                Inches(0.3), Inches(0.3), size=Pt(12), bold=True,
                color=C_BG, align=PP_ALIGN.CENTER)
            txt(s, item, bx + Inches(0.65), iy + Inches(0.04),
                cw - Inches(0.78), Inches(0.42), size=Pt(11), color=C_LGRAY)

    txt(s, ">", Inches(6.2), Inches(3.75), Inches(0.7), Inches(0.7),
        size=Pt(36), bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    # Metrics row — real numbers from evaluation sprint
    metrics = [("180",  "Benchmark\nSpecifications", C_GREEN),
               ("20",   "Executable\nContracts",     C_CYAN),
               ("200",  "Adversarial\nScenarios",    C_YELLOW),
               ("33",   "Replay\nRegression Cases",  C_GREEN_L)]
    for mi, (val, lbl, mc) in enumerate(metrics):
        mx = Inches(0.55) + mi * Inches(3.1)
        my = Inches(6.62)
        rect(s, mx, my, Inches(2.75), Inches(0.72), C_CARD2)
        txt(s, val, mx + Inches(0.15), my + Inches(0.04),
            Inches(0.8), Inches(0.55), size=Pt(26), bold=True, color=mc)
        txt(s, lbl, mx + Inches(1.0), my + Inches(0.08),
            Inches(1.6), Inches(0.5), size=Pt(10), color=C_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE 9 — EVALUATION INFRASTRUCTURE (new)
# ─────────────────────────────────────────────────────────
def slide_eval_infra(prs):
    s = blank(prs); bg(s)
    slide_num(s, 9); tag(s, "EVALUATION INFRASTRUCTURE")
    title(s, "Evaluation Infrastructure — Built from Scratch")
    subtitle_line(s, "Weeks of dedicated engineering that underpins every readiness claim in this deck")
    divider(s, Inches(1.65))

    registries = [
        ("Benchmark Registry",  C_GREEN,  "180",
         "Contract Scenarios",
         "One specification per pattern variant. Positive, negative, and edge-case paths. Versioned and diffable."),
        ("Replay Corpus",       C_YELLOW, "33",
         "Historical Regressions",
         "Every past failure that reached production is now a permanent regression case. No bug re-enters undetected."),
        ("Adversarial Registry",C_ORANGE, "200",
         "Attack Scenarios",
         "Reentrancy, overflow, authority theft, token inflation, category spoofing — all systematically encoded."),
        ("Real-World Corpus",   C_CYAN,   "28",
         "BCH Contracts",
         "Live contracts from the BCH ecosystem used as ground-truth evaluation targets. Validates real coverage."),
    ]

    rw = Inches(5.85)
    rh = Inches(1.45)
    for ri, (name, clr, count, unit, desc) in enumerate(registries):
        rx = Inches(0.55) + (ri % 2) * (rw + Inches(0.28))
        ry = Inches(1.82) + (ri // 2) * (rh + Inches(0.2))
        rect(s, rx, ry, rw, rh, C_CARD)
        rect(s, rx, ry, Inches(0.055), rh, clr)

        # Count badge
        txt(s, count, rx + Inches(0.18), ry + Inches(0.08),
            Inches(0.95), Inches(0.72), size=Pt(38), bold=True, color=clr)
        txt(s, unit, rx + Inches(1.18), ry + Inches(0.32),
            Inches(1.5), Inches(0.35), size=Pt(10), color=C_GRAY)

        # Name + desc
        txt(s, name, rx + Inches(2.85), ry + Inches(0.08),
            rw - Inches(3.0), Inches(0.38), size=Pt(13), bold=True, color=C_WHITE)
        txt(s, desc, rx + Inches(2.85), ry + Inches(0.55),
            rw - Inches(3.0), Inches(0.75), size=Pt(10), color=C_GRAY)

    # CI Validation strip
    rect(s, Inches(0.55), Inches(5.62), W - Inches(1.1), Inches(0.72), C_CARD2)
    rect(s, Inches(0.55), Inches(5.62), Inches(0.055), Inches(0.72), C_GREEN_L)
    txt(s, "CI Validation", Inches(0.72), Inches(5.68),
        Inches(2.2), Inches(0.35), size=Pt(13), bold=True, color=C_WHITE)
    txt(s, "Zero-cost deterministic evaluation — every benchmark runs on each code change. No environment drift. No human in the loop.",
        Inches(3.1), Inches(5.68), Inches(9.65), Inches(0.52),
        size=Pt(11), color=C_LGRAY)

    rect(s, Inches(0.55), Inches(6.5), W - Inches(1.1), Inches(0.38), C_CARD2)
    txt(s, "This infrastructure did not exist at Hackcelerator start. It was built in response to measurement failures encountered during development.",
        Inches(0.72), Inches(6.53), W - Inches(1.3), Inches(0.32),
        size=Pt(9), color=C_GRAY, italic=True)


# ─────────────────────────────────────────────────────────
# SLIDE 10 — SPLIT PAYMENT BREAKTHROUGH (hero story format)
# ─────────────────────────────────────────────────────────
def slide_08(prs):
    s = blank(prs); bg(s)
    slide_num(s, 10); tag(s, "CASE STUDY")
    title(s, "Split Payment: From 0% to 100%")
    subtitle_line(s, "How benchmarking exposed hidden blockers and drove convergence")
    divider(s, Inches(1.65))

    # Before / After progress bars
    for label, score, clr, y in [
        ("BEFORE — Initial compile score",      0,   C_ORANGE, Inches(1.85)),
        ("AFTER — Validation subset score",     100, C_GREEN,  Inches(2.62)),
    ]:
        txt(s, label, Inches(0.55), y, Inches(7), Inches(0.3),
            size=Pt(9), bold=True, color=clr)
        rect(s, Inches(0.55), y + Inches(0.35), Inches(10.5), Inches(0.4), C_CARD)
        if score > 0:
            rect(s, Inches(0.55), y + Inches(0.35),
                 Inches(10.5 * score / 100), Inches(0.4), clr)
        txt(s, f"{score}%", Inches(11.15), y + Inches(0.38), Inches(0.95), Inches(0.35),
            size=Pt(13), bold=True, color=clr)

    # Story cards
    cards = [
        ("PROBLEM",    "Compile score: 0%\nContract generated but\nfailed every evaluation",    C_ORANGE),
        ("ROOT CAUSE", "Routing bug in generator:\nN-output path assumption\n+ structural false positive in audit", C_YELLOW),
        ("FIX APPLIED","N-output routing corrected.\nAudit structural check\naligned to compile phase",  C_CYAN),
        ("RESULT",     "100% convergence\non validation subset.\nBenchmark gap exposed the real blocker.", C_GREEN),
    ]

    cw = Inches(2.95)
    for ci, (stage, desc, clr) in enumerate(cards):
        cx = Inches(0.55) + ci * (cw + Inches(0.22))
        cy = Inches(3.35)
        ch = Inches(2.65)
        rect(s, cx, cy, cw, ch, C_CARD)
        rect(s, cx, cy, cw, Inches(0.45), C_CARD2)
        rect(s, cx, cy, Inches(0.055), ch, clr)
        txt(s, stage, cx + Inches(0.18), cy + Inches(0.08),
            cw - Inches(0.25), Inches(0.35), size=Pt(10), bold=True, color=clr)
        txt(s, desc, cx + Inches(0.18), cy + Inches(0.55),
            cw - Inches(0.25), Inches(1.95), size=Pt(11), color=C_LGRAY)

        if ci < 3:
            txt(s, ">", cx + cw + Inches(0.02), cy + Inches(0.95),
                Inches(0.22), Inches(0.4), size=Pt(14), color=C_CARD2,
                align=PP_ALIGN.CENTER)

    # Lesson
    rect(s, Inches(0.55), Inches(6.2), W - Inches(1.1), Inches(0.65), C_CARD2)
    rect(s, Inches(0.55), Inches(6.2), Inches(0.055), Inches(0.65), C_GREEN)
    txt(s,
        "KEY LESSON:  Benchmarking does not just measure progress — it exposes hidden blockers that manual review misses entirely.",
        Inches(0.72), Inches(6.26), W - Inches(1.3), Inches(0.52),
        size=Pt(11), bold=True, color=C_WHITE)


# ─────────────────────────────────────────────────────────
# SLIDE 9 — PATTERN STABILIZATION (story format, accurate RCAs)
# ─────────────────────────────────────────────────────────
def slide_09(prs):
    s = blank(prs); bg(s)
    slide_num(s, 11); tag(s, "PATTERN STABILIZATION")
    title(s, "Escrow · Multisig · Timelock — How They Stabilized")
    subtitle_line(s, "Same root-cause analysis framework applied across all three patterns")
    divider(s, Inches(1.65))

    stories = [
        ("ESCROW", C_GREEN, [
            ("PROBLEM",    "Low benchmark scores despite\nvalid contract generation"),
            ("ROOT CAUSE", "Evaluator / benchmark mismatch:\nexpected output format did not\nalign with generator output"),
            ("FIX",        "Aligned benchmark evaluator\nto match generator contract\nstructure and output schema"),
            ("RESULT",     "Score: 1.0\nAcross all positive\nbenchmark paths"),
        ]),
        ("MULTISIG", C_CYAN, [
            ("PROBLEM",    "Scoring inconsistent across\ndifferent M-of-N configurations"),
            ("ROOT CAUSE", "Feature detection mismatch:\ntoken_validation field mapping\nnot correctly extracted"),
            ("FIX",        "Updated feature mapping layer\nto correctly parse multisig\ncapability declarations"),
            ("RESULT",     "Score: 1.0\nAll threshold configurations\npass evaluation"),
        ]),
        ("TIMELOCK", C_YELLOW, [
            ("PROBLEM",    "0.093 average score despite\ncorrect timelock logic"),
            ("ROOT CAUSE", "Evaluator coverage gaps:\nmissing evaluator mappings for\nabsolute + relative locktime cases"),
            ("FIX",        "Added missing evaluator mappings\nfor all locktime variants in\nbenchmark suite"),
            ("RESULT",     "Score: 1.0\nAbsolute & relative\ntimelocks both covered"),
        ]),
    ]

    pw = Inches(3.85)
    for pi, (name, clr, steps) in enumerate(stories):
        px = Inches(0.55) + pi * (pw + Inches(0.22))
        # Header
        rect(s, px, Inches(1.82), pw, Inches(0.5), C_CARD2)
        rect(s, px, Inches(1.82), Inches(0.055), Inches(0.5), clr)
        txt(s, name, px + Inches(0.18), Inches(1.87),
            pw - Inches(0.25), Inches(0.38), size=Pt(14), bold=True, color=C_WHITE)

        for si, (stage, desc) in enumerate(steps):
            sy = Inches(2.42) + si * Inches(1.12)
            sh = Inches(1.0)
            is_result = stage == "RESULT"
            card_bg = C_CARD2 if is_result else C_CARD
            rect(s, px, sy, pw, sh, card_bg)
            rect(s, px, sy, Inches(0.055), sh, clr if is_result else C_CARD2)
            txt(s, stage, px + Inches(0.15), sy + Inches(0.06),
                Inches(1.2), Inches(0.28), size=Pt(8), bold=True,
                color=clr if is_result else C_GRAY)
            txt(s, desc, px + Inches(0.15), sy + Inches(0.35),
                pw - Inches(0.22), Inches(0.58), size=Pt(10),
                color=C_WHITE if is_result else C_LGRAY)

            if si < 3:
                txt(s, "v", px + pw / 2 - Inches(0.15), sy + sh,
                    Inches(0.3), Inches(0.18), size=Pt(9),
                    color=C_CARD2, align=PP_ALIGN.CENTER)

    rect(s, Inches(0.55), Inches(6.85), W - Inches(1.1), Inches(0.32), C_CARD2)
    txt(s, "Pattern: Identify measurement gaps first. Fixing the evaluator is often faster than fixing the generator.",
        Inches(0.72), Inches(6.88), W - Inches(1.3), Inches(0.28),
        size=Pt(9), color=C_GRAY, italic=True)


# ─────────────────────────────────────────────────────────
# SLIDE 10 — VAULT RESEARCH (revised claim)
# ─────────────────────────────────────────────────────────
def slide_10(prs):
    s = blank(prs); bg(s)
    slide_num(s, 12); tag(s, "VAULT RESEARCH")
    title(s, "Vault — Emerging as the Most Advanced Pattern Family")
    subtitle_line(s, "Six distinct vault architectures under active development")
    divider(s, Inches(1.65))

    # Hero left panel
    rect(s, Inches(0.55), Inches(1.82), Inches(4.2), Inches(5.12), C_CARD)
    rect(s, Inches(0.55), Inches(1.82), Inches(0.055), Inches(5.12), C_YELLOW)
    txt(s, "VAULT\nARCHITECTURE", Inches(0.72), Inches(1.92),
        Inches(3.85), Inches(0.85), size=Pt(22), bold=True, color=C_WHITE)
    txt(s, "Combines covenants, timelocks,\nmultisig, and CashTokens into\ncomplex multi-path custody systems.",
        Inches(0.72), Inches(2.88), Inches(3.85), Inches(0.85),
        size=Pt(11), color=C_GRAY)

    # Status
    for label, val, clr, y in [
        ("Generation",  "67% — active improvement",  C_YELLOW,  Inches(3.88)),
        ("Audit",       "Target: 95% after timeout fixes", C_GREEN_L, Inches(4.42)),
        ("Status",      "Under active development",  C_CYAN,    Inches(4.96)),
    ]:
        rect(s, Inches(0.72), y, Inches(0.1), Inches(0.32), clr)
        txt(s, f"{label}:  {val}", Inches(0.96), y - Inches(0.02),
            Inches(3.5), Inches(0.38), size=Pt(11), color=C_LGRAY)

    rect(s, Inches(0.72), Inches(5.6), Inches(3.7), Inches(1.05), C_CARD2)
    txt(s,
        "Vault is emerging as the most advanced\npattern family in NexOps — not yet\ncomplete but the clearest frontier.",
        Inches(0.88), Inches(5.68), Inches(3.4), Inches(0.88),
        size=Pt(10), color=C_WHITE, italic=True)

    # Right — 6 vault types as list
    vaults = [
        ("Recovery Vault",       C_GREEN,   "Time-locked funds with emergency recovery path"),
        ("Delayed Withdrawals",  C_GREEN_L, "Configurable delay window with cancellation"),
        ("Emergency Paths",      C_YELLOW,  "Multi-sig bypass for urgent access scenarios"),
        ("Tiered Vaults",        C_YELLOW,  "Multiple spending tiers with distinct auth levels"),
        ("Treasury Protection",  C_CYAN,    "DAO-grade multi-party custody with governance"),
        ("Covenant Vaults",      C_ORANGE,  "Self-enforcing spend policies, no oracle required"),
    ]

    rx = Inches(5.05)
    for vi, (name, clr, desc) in enumerate(vaults):
        vy = Inches(1.82) + vi * Inches(0.85)
        rect(s, rx, vy, Inches(7.73), Inches(0.75), C_CARD)
        rect(s, rx, vy, Inches(0.055), Inches(0.75), clr)
        txt(s, name, rx + Inches(0.18), vy + Inches(0.06),
            Inches(2.5), Inches(0.38), size=Pt(12), bold=True, color=C_WHITE)
        txt(s, desc, rx + Inches(2.8), vy + Inches(0.08),
            Inches(4.75), Inches(0.55), size=Pt(10), color=C_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE 11 — COMPOSITION RESEARCH
# ─────────────────────────────────────────────────────────
def slide_11(prs):
    s = blank(prs); bg(s)
    slide_num(s, 13); tag(s, "COMPOSITION RESEARCH")
    title(s, "Beyond Single Contracts")
    subtitle_line(s, "Future on-chain systems are compositions of proven security primitives")
    divider(s, Inches(1.65))

    combos = [
        ("Multisig + Timelock",   C_GREEN,   "DAO governance with time-bounded proposals"),
        ("Escrow + Hashlock",     C_CYAN,    "Atomic cross-chain swap primitive"),
        ("Vault + Timelock",      C_YELLOW,  "Treasury with delayed withdrawal windows"),
        ("Payroll + Governance",  C_GREEN_L, "FT-denominated automated payroll DAO"),
        ("FT + NFT Treasury",     C_ORANGE,  "Hybrid token vault with NFT access keys"),
        ("Escrow + Multisig",     C_GREEN,   "Multi-party arbitration escrow system"),
    ]

    cw = Inches(4.05)
    ch = Inches(1.55)
    for ci, (name, clr, desc) in enumerate(combos):
        cx = Inches(0.55) + (ci % 3) * (cw + Inches(0.2))
        cy = Inches(1.82) + (ci // 3) * (ch + Inches(0.2))
        rect(s, cx, cy, cw, ch, C_CARD)
        rect(s, cx, cy, Inches(0.055), ch, clr)
        # Pills
        parts = name.split(" + ")
        px = cx + Inches(0.18)
        for part in parts:
            pw2 = Inches(1.52)
            rect(s, px, cy + Inches(0.14), pw2, Inches(0.3), C_CARD2)
            txt(s, part, px + Inches(0.05), cy + Inches(0.14), pw2, Inches(0.28),
                size=Pt(9), bold=True, color=clr, align=PP_ALIGN.CENTER)
            px += pw2 + Inches(0.08)
            if part != parts[-1]:
                txt(s, "+", px - Inches(0.1), cy + Inches(0.14), Inches(0.15), Inches(0.28),
                    size=Pt(9), color=C_GRAY, align=PP_ALIGN.CENTER)

        txt(s, desc, cx + Inches(0.18), cy + Inches(0.6),
            cw - Inches(0.25), Inches(0.72), size=Pt(11), color=C_LGRAY)

    rect(s, Inches(0.55), Inches(6.58), W - Inches(1.1), Inches(0.52), C_CARD2)
    rect(s, Inches(0.55), Inches(6.58), Inches(0.055), Inches(0.52), C_GREEN)
    txt(s,
        "Composition Engine (next milestone): automatically verify security invariants across composed contract systems.",
        Inches(0.72), Inches(6.63), W - Inches(1.3), Inches(0.42),
        size=Pt(11), bold=True, color=C_WHITE)


# ─────────────────────────────────────────────────────────
# SLIDE 12 — CURRENT STATE (no invented percentages)
# ─────────────────────────────────────────────────────────
def slide_12(prs):
    s = blank(prs); bg(s)
    slide_num(s, 14); tag(s, "PLATFORM TODAY")
    title(s, "NexOps Platform — Current Capabilities")
    subtitle_line(s, "Honest capability map — what is production-ready vs what is validated")
    divider(s, Inches(1.65))

    caps = [
        ("Contract Generation",    C_GREEN,
         "Production-ready on core BCH primitives",
         "9 pattern families. Full CashScript output. Stable for core cases."),
        ("Security Auditing",      C_GREEN_L,
         "Semantic + structural verification",
         "Intent invariant engine + compile verification. Multi-layer pipeline."),
        ("Eval Infrastructure",    C_CYAN,
         "Validated benchmark infrastructure",
         "180 specs · 20 contracts · 200 adversarial · 33 regression cases."),
        ("CashTokens Support",     C_GREEN,
         "FT + NFT full support implemented",
         "Fungible, non-fungible, authority, hybrid covenant patterns."),
        ("Composition Engine",     C_ORANGE,
         "Research phase — next major milestone",
         "Composition engine in design. Multi-contract audit reasoning planned."),
    ]

    bx = Inches(0.55)
    bw = Inches(8.8)

    for i, (cap, clr, status, detail) in enumerate(caps):
        by = Inches(1.82) + i * Inches(0.95)
        bh = Inches(0.82)
        rect(s, bx, by, bw, bh, C_CARD)
        rect(s, bx, by, Inches(0.055), bh, clr)

        txt(s, cap, bx + Inches(0.18), by + Inches(0.06),
            Inches(2.4), Inches(0.38), size=Pt(12), bold=True, color=C_WHITE)
        txt(s, status, bx + Inches(2.7), by + Inches(0.06),
            Inches(2.8), Inches(0.38), size=Pt(11), bold=True, color=clr)
        txt(s, detail, bx + Inches(5.65), by + Inches(0.06),
            Inches(3.0), Inches(0.65), size=Pt(9), color=C_GRAY)

    # Right summary
    rx = Inches(10.0)
    rect(s, rx, Inches(1.82), Inches(2.78), Inches(4.72), C_CARD)
    rect(s, rx, Inches(1.82), Inches(0.055), Inches(4.72), C_GREEN)
    txt(s, "MATURITY\nSUMMARY", rx + Inches(0.18), Inches(1.9),
        Inches(2.4), Inches(0.7), size=Pt(14), bold=True, color=C_WHITE)

    tiers_r = [
        ("Production",  ["Generation", "CashTokens"],            C_GREEN),
        ("Validated",   ["Auditing", "Eval Infrastructure"],     C_CYAN),
        ("Research",    ["Composition Engine"],                   C_ORANGE),
    ]
    ry = Inches(2.75)
    for tier_n, tier_items, tc in tiers_r:
        oval(s, rx + Inches(0.18), ry + Inches(0.06), Inches(0.2), Inches(0.2), tc)
        txt(s, tier_n, rx + Inches(0.52), ry, Inches(2.1), Inches(0.3),
            size=Pt(10), bold=True, color=tc)
        for item in tier_items:
            ry += Inches(0.32)
            txt(s, item, rx + Inches(0.52), ry, Inches(2.1), Inches(0.28),
                size=Pt(10), color=C_LGRAY)
        ry += Inches(0.45)

    rect(s, Inches(0.55), Inches(6.72), W - Inches(1.1), Inches(0.32), C_CARD2)
    txt(s, "All status labels derived from benchmark convergence data, not estimated percentages.",
        Inches(0.72), Inches(6.75), W - Inches(1.3), Inches(0.28),
        size=Pt(9), color=C_GRAY, italic=True)


# ─────────────────────────────────────────────────────────
# SLIDE 13 — WHAT WE LEARNED (new slide)
# ─────────────────────────────────────────────────────────
def slide_13(prs):
    s = blank(prs); bg(s)
    slide_num(s, 15); tag(s, "ENGINEERING LESSONS")
    title(s, "What We Learned")
    subtitle_line(s, "Key insights that changed how we approach security infrastructure")
    divider(s, Inches(1.65))

    lessons = [
        ("Benchmarking reveals hidden failures",
         C_GREEN,
         "Manual review passes contracts that benchmarks catch. We found real blockers only after building repeatable suites. You cannot trust what you cannot measure."),
        ("Measurement bugs look like generation bugs",
         C_YELLOW,
         "Escrow, Multisig, and Timelock all appeared broken. The generator was correct. The evaluator was wrong. Always audit the test before auditing the contract."),
        ("Security validation must understand intent",
         C_CYAN,
         "Structural checks are necessary but not sufficient. A contract can be syntactically valid and semantically dangerous. Intent invariants are the missing layer."),
        ("Composition is harder than individual patterns",
         C_ORANGE,
         "Each pattern stabilizes in isolation. Composing two stable patterns creates new security surface. Cross-pattern invariants require entirely different reasoning."),
        ("CashTokens require separate semantic reasoning",
         C_GREEN_L,
         "Token logic and script logic interact in non-obvious ways. Category continuity, amount preservation, and authority semantics need dedicated audit layers."),
    ]

    lw = Inches(12.23)
    lh = Inches(0.9)
    for i, (title_l, clr, body) in enumerate(lessons):
        ly = Inches(1.82) + i * (lh + Inches(0.1))
        rect(s, Inches(0.55), ly, lw, lh, C_CARD)
        rect(s, Inches(0.55), ly, Inches(0.055), lh, clr)

        # Number
        txt(s, f"{i+1:02d}", Inches(0.72), ly + Inches(0.28),
            Inches(0.45), Inches(0.35), size=Pt(13), bold=True, color=clr)
        # Title
        txt(s, title_l, Inches(1.22), ly + Inches(0.06),
            Inches(3.9), Inches(0.38), size=Pt(12), bold=True, color=C_WHITE)
        # Body
        txt(s, body, Inches(5.2), ly + Inches(0.06),
            Inches(7.4), Inches(0.75), size=Pt(10), color=C_GRAY)


# ─────────────────────────────────────────────────────────
# SLIDE 14 — ROADMAP (toned-down, composition engine first)
# ─────────────────────────────────────────────────────────
def slide_14(prs):
    s = blank(prs); bg(s)
    slide_num(s, 16); tag(s, "ROADMAP")
    title(s, "Next Phase — Building the Intelligence Layer")
    subtitle_line(s, "Milestones grounded in work already completed")
    divider(s, Inches(1.65))

    phases = [
        ("NEAR TERM\nQ3 – Q4 2026", C_GREEN, [
            ("Pattern Composition Engine",     "First milestone — directly follows completed primitive work"),
            ("Composite Benchmark Suites",     "Extend existing suites to cover multi-contract scenarios"),
            ("Multi-Contract Audit Reasoning", "Security invariants across composed systems"),
            ("Production CashTokens Coverage", "Complete FT + NFT + hybrid covenant benchmark parity"),
        ]),
        ("MID TERM\nQ1 – Q2 2027", C_YELLOW, [
            ("DAO Treasury Systems",           "Multi-sig + governance + token custody patterns"),
            ("Oracle-Aware Contracts",         "External data feed integration with security model"),
            ("Governance Frameworks",          "On-chain voting with composable contract primitives"),
            ("Cross-Pattern Security Proofs",  "Formal invariant verification across pattern families"),
        ]),
        ("LONG TERM\n2027+", C_GRAY, [
            ("Autonomous Review Layer",        "Pipeline-integrated contract review on generation"),
            ("Enterprise Audit API",           "Programmatic access for dApp developers"),
            ("BCH Ecosystem Standard",         "Open benchmark suite adopted across BCH tooling"),
            ("Extended Chain Support",         "If warranted by ecosystem demand"),
        ]),
    ]

    pw = Inches(3.85)
    for pi, (phase_title, clr, items) in enumerate(phases):
        px = Inches(0.55) + pi * (pw + Inches(0.22))

        rect(s, px, Inches(1.82), pw, Inches(0.65), C_CARD2)
        rect(s, px, Inches(1.82), Inches(0.055), Inches(0.65), clr)
        txt(s, phase_title, px + Inches(0.18), Inches(1.86),
            pw - Inches(0.25), Inches(0.55), size=Pt(11), bold=True, color=C_WHITE)

        for ii, (item_title, item_desc) in enumerate(items):
            iy = Inches(2.58) + ii * Inches(1.0)
            rect(s, px, iy, pw, Inches(0.88), C_CARD)
            rect(s, px, iy, Inches(0.055), Inches(0.88), clr)
            txt(s, item_title, px + Inches(0.18), iy + Inches(0.04),
                pw - Inches(0.25), Inches(0.38), size=Pt(11), bold=True, color=C_WHITE)
            txt(s, item_desc, px + Inches(0.18), iy + Inches(0.45),
                pw - Inches(0.25), Inches(0.35), size=Pt(9), color=C_GRAY)

    # Timeline
    rect(s, Inches(0.55), Inches(6.72), Inches(12.23), Inches(0.04), C_CARD2)
    for ti, (lbl, clr) in enumerate([("Now", C_GREEN), ("Q3 '26", C_GREEN_L),
                                      ("Q1 '27", C_YELLOW), ("2027+", C_GRAY)]):
        tx = Inches(0.55) + ti * Inches(4.08)
        oval(s, tx - Inches(0.1), Inches(6.62), Inches(0.22), Inches(0.22), clr)
        txt(s, lbl, tx - Inches(0.4), Inches(6.88), Inches(0.9), Inches(0.28),
            size=Pt(9), color=clr, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDE 15 — KEY ACHIEVEMENTS (geometric indicators, no emoji)
# ─────────────────────────────────────────────────────────
def slide_15(prs):
    s = blank(prs); bg(s)
    slide_num(s, 17); tag(s, "KEY ACHIEVEMENTS")
    title(s, "What We Built — Results That Matter")
    subtitle_line(s, "Six major milestones delivered across security, benchmarking, and token infrastructure")
    divider(s, Inches(1.65))

    achievements = [
        ("01", "BCH Hackcelerator\nSelection",        C_YELLOW,
         "Selected from global applicant pool.\nValidated concept at program start."),
        ("02", "CashTokens\nFull Implementation",     C_CYAN,
         "FT + NFT + authority + hybrid covenant\npatterns all implemented."),
        ("03", "Benchmark\nInfrastructure",           C_GREEN,
         "5+ suites, 100+ cases. Repeatable,\nversioned, adversarial coverage."),
        ("04", "Semantic\nAudit Engine",               C_GREEN_L,
         "Intent invariant pipeline. Structural\n+ semantic + compile layers."),
        ("05", "Pattern Stabilization\n(3 families)",  C_YELLOW,
         "Escrow, Multisig, Timelock all converge\nto 1.0 on positive benchmark paths."),
        ("06", "Vault Architecture\nResearch",         C_ORANGE,
         "6 advanced vault patterns.\nCovenants + multisig + timelocks."),
    ]

    aw = Inches(3.9)
    ah = Inches(2.05)
    for ai, (num, ttl, clr, desc) in enumerate(achievements):
        ax = Inches(0.55) + (ai % 3) * (aw + Inches(0.22))
        ay = Inches(1.82) + (ai // 3) * (ah + Inches(0.2))
        rect(s, ax, ay, aw, ah, C_CARD)
        rect(s, ax, ay, aw, Inches(0.055), clr)

        # Number badge
        rect(s, ax + Inches(0.18), ay + Inches(0.18), Inches(0.46), Inches(0.46), C_CARD2)
        txt(s, num, ax + Inches(0.18), ay + Inches(0.2), Inches(0.46), Inches(0.38),
            size=Pt(12), bold=True, color=clr, align=PP_ALIGN.CENTER)

        txt(s, ttl, ax + Inches(0.78), ay + Inches(0.12),
            aw - Inches(0.95), Inches(0.65), size=Pt(13), bold=True, color=C_WHITE)
        txt(s, desc, ax + Inches(0.18), ay + Inches(0.88),
            aw - Inches(0.3), Inches(0.92), size=Pt(10), color=C_GRAY)

    rect(s, Inches(0.55), Inches(6.68), W - Inches(1.1), Inches(0.42), C_CARD2)
    rect(s, Inches(0.55), Inches(6.68), Inches(0.055), Inches(0.42), C_GREEN)
    txt(s, "Security-first architecture throughout  •  BCH ecosystem leader  •  Mentor-ready documentation",
        Inches(0.72), Inches(6.72), W - Inches(1.3), Inches(0.32),
        size=Pt(10), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────
# SLIDE 16 — VISION
# ─────────────────────────────────────────────────────────
def slide_16(prs):
    s = blank(prs); bg(s)

    # Dot grid
    for row in range(9):
        for col in range(22):
            x = Inches(0.12 + col * 0.62)
            y = Inches(0.2 + row * 0.85 + (0.42 if col % 2 else 0))
            oval(s, x, y, Inches(0.045), Inches(0.045),
                 RGBColor(0x1a, 0x26, 0x35))

    # Green left spine
    rect(s, Inches(0.55), Inches(0.55), Inches(0.055), Inches(4.2), C_GREEN)

    # Vision statement
    txt(s,
        "NexOps aims to become the\nsecurity layer between human\nintent and on-chain execution.",
        Inches(0.72), Inches(0.6), Inches(8.5), Inches(2.1),
        size=Pt(34), bold=True, color=C_WHITE)

    # Flow diagram
    flow = [
        ("Human\nIntent",        C_LGRAY,  Inches(2.1)),
        (">",                    C_CARD2,  Inches(0.5)),
        ("NexOps\nProtocol",     C_GREEN,  Inches(2.1)),
        (">",                    C_CARD2,  Inches(0.5)),
        ("Secure Smart\nContracts",C_CYAN, Inches(2.4)),
        (">",                    C_CARD2,  Inches(0.5)),
        ("On-Chain\nState",      C_YELLOW, Inches(2.1)),
    ]
    fx = Inches(0.72)
    for text, clr, fw in flow:
        is_arrow = text == ">"
        fh = Inches(1.0)
        fy = Inches(3.0)
        if not is_arrow:
            rect(s, fx, fy, fw, fh, C_CARD)
            rect(s, fx, fy, Inches(0.055), fh, clr)
        txt(s, text, fx, fy + (Inches(0.2) if not is_arrow else Inches(0.32)),
            fw, Inches(0.55),
            size=Pt(12 if not is_arrow else 20),
            bold=(not is_arrow), color=clr,
            align=PP_ALIGN.CENTER if not is_arrow else PP_ALIGN.CENTER)
        fx += fw + Inches(0.06)

    # Four pillars
    pillars = [("Generate.", C_GREEN), ("Verify.", C_GREEN_L),
               ("Benchmark.", C_YELLOW), ("Secure.", C_ORANGE)]
    pw2 = Inches(2.78)
    for pi, (text, clr) in enumerate(pillars):
        px = Inches(0.72) + pi * (pw2 + Inches(0.25))
        rect(s, px, Inches(4.3), pw2, Inches(0.6), C_CARD)
        rect(s, px, Inches(4.88), pw2, Inches(0.04), clr)
        txt(s, text, px, Inches(4.32), pw2, Inches(0.52),
            size=Pt(17), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # Footer
    rect(s, 0, Inches(5.25), W, Inches(0.055), C_GREEN)
    rect(s, 0, Inches(5.3), W, Inches(2.2), RGBColor(0x04, 0x09, 0x0f))

    txt(s, "NEXOPS PROTOCOL",
        Inches(0.72), Inches(5.55), Inches(6), Inches(0.75),
        size=Pt(36), bold=True, color=C_WHITE)
    txt(s, "Security-First Smart Contract Infrastructure for Bitcoin Cash",
        Inches(0.72), Inches(6.3), Inches(8), Inches(0.42),
        size=Pt(13), color=C_GRAY)

    rect(s, Inches(10.5), Inches(5.5), Inches(2.3), Inches(1.05), C_CARD2, C_GREEN, Pt(1.2))
    txt(s, "BCH\nHackcelerator", Inches(10.58), Inches(5.58),
        Inches(2.15), Inches(0.88), size=Pt(13), bold=True,
        color=C_GREEN, align=PP_ALIGN.CENTER)

    bottom_bar(s)


# ─────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_01(prs)              # 01 Cover
    slide_02(prs)              # 02 Where We Started
    slide_03(prs)              # 03 Evolution Timeline
    slide_04(prs)              # 04 Pattern Coverage Matrix
    slide_05(prs)              # 05 CashTokens (+ benchmark proof points)
    slide_06(prs)              # 06 Security Engine
    slide_audit_evolution(prs) # 07 Security Audit Evolution  [NEW]
    slide_07(prs)              # 08 Benchmark System (real numbers)
    slide_eval_infra(prs)      # 09 Evaluation Infrastructure [NEW]
    slide_08(prs)              # 10 Split Payment Case Study
    slide_09(prs)              # 11 Pattern Stabilization
    slide_10(prs)              # 12 Vault Research
    slide_11(prs)              # 13 Composition Research
    slide_12(prs)              # 14 Platform Today (updated maturity)
    slide_13(prs)              # 15 What We Learned
    slide_14(prs)              # 16 Roadmap
    slide_15(prs)              # 17 Key Achievements
    slide_16(prs)              # 18 Vision

    out = r"d:\downloadds\nexmcp\nexops-mcp\docs\NexOps_Protocol_Deck_v3.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
